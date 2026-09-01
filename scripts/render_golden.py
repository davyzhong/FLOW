"""Freeze, approve, and render the golden publishing fixtures.

Used by scripts/test_publishing_golden.sh: renders pptx/xlsx/html from one
frozen Report Snapshot, extracts the canonical key values, and writes them
next to the artifacts for cross-format verification.
"""

from __future__ import annotations

import argparse
import json
import os
import sys
from io import BytesIO
from pathlib import Path

REPOSITORY_ROOT = Path(__file__).resolve().parents[1]
API_ROOT = REPOSITORY_ROOT / "services" / "api"
PUBLISHING_TESTS = API_ROOT / "tests" / "publishing"
INTEGRATION_TESTS = API_ROOT / "tests" / "integration"
for entry in (
    str(API_ROOT / "src"),
    str(PUBLISHING_TESTS),
    str(INTEGRATION_TESTS),
    str(API_ROOT / "tests"),
    str(API_ROOT),
):
    if entry not in sys.path:
        sys.path.insert(0, entry)

from openpyxl import load_workbook  # noqa: E402
from pptx import Presentation  # noqa: E402

from flow_api.infrastructure.db import get_session_factory  # noqa: E402


def _ensure_env() -> None:
    os.environ.setdefault(
        "DATABASE_URL", "postgresql+psycopg://flow:flow_dev_only@127.0.0.1:5432/flow"
    )
    os.environ.setdefault("REDIS_URL", "redis://127.0.0.1:6379/0")
    os.environ.setdefault("S3_ENDPOINT_URL", "http://127.0.0.1:9000")
    os.environ.setdefault("S3_BUCKET", "flow")
    os.environ.setdefault("S3_ACCESS_KEY", "flow")
    os.environ.setdefault("S3_SECRET_KEY", "flow_dev_only")


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--out", required=True)
    arguments = parser.parse_args()
    out_dir = Path(arguments.out)
    out_dir.mkdir(parents=True, exist_ok=True)
    _ensure_env()

    from publishing_support import fresh_approved_report  # type: ignore[import-not-found]

    with get_session_factory()() as session:
        report, view = fresh_approved_report(session)

    (out_dir / "report.pptx").write_bytes(render_pptx(view))
    (out_dir / "report.xlsx").write_bytes(render_xlsx(view))
    (out_dir / "report.html").write_bytes(render_html(view))

    key_values = view.key_values()
    key_values["title"] = view.identity.title

    presentation = Presentation(BytesIO((out_dir / "report.pptx").read_bytes()))
    pptx_text = "\n".join(
        shape.text_frame.text
        for slide in presentation.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    workbook = load_workbook(BytesIO((out_dir / "report.xlsx").read_bytes()))
    xlsx_values = {
        str(cell.value)
        for sheet in workbook.worksheets
        for row in sheet.iter_rows()
        for cell in row
        if cell.value is not None
    }
    html_text = (out_dir / "report.html").read_text(encoding="utf-8")

    missing: list[str] = []
    for key, value in key_values.items():
        if value in (None, ""):
            continue
        if value not in pptx_text:
            missing.append(f"pptx:{key}")
        if value not in xlsx_values:
            missing.append(f"xlsx:{key}")
        if value not in html_text:
            missing.append(f"html:{key}")
    (out_dir / "key_values.json").write_text(
        json.dumps({"key_values": key_values, "missing": missing}, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    if missing:
        print(f"golden cross-format check FAILED: missing {missing}")
        return 1
    print(
        f"golden cross-format check passed (pptx/xlsx/html): {len(key_values)} key values."
    )
    return 0


from flow_api.publishing.renderers import render_html, render_pptx, render_xlsx  # noqa: E402


if __name__ == "__main__":
    raise SystemExit(main())
