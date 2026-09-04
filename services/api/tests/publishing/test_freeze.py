from __future__ import annotations

from io import BytesIO

import pytest
from integration.analysis_run_support import (
    analysis_session_fixture as analysis_session,  # noqa: F401
)
from integration.intake_service_support import (
    intake_session_fixture as intake_session_fixture,  # noqa: F401
)
from integration.metric_snapshot_support import (
    metric_session_fixture as metric_session_fixture,  # noqa: F401
)
from openpyxl import load_workbook
from pptx import Presentation
from publishing.publishing_support import publish_analysis_run
from publishing.publishing_support import (
    publishing_session_fixture as _publishing_session_fixture,  # noqa: F401
)
from sqlalchemy import select
from sqlalchemy.orm import Session

from flow_api.infrastructure.models.analytics import Conclusion, Evidence, Finding
from flow_api.infrastructure.models.publishing import PublicationAttempt
from flow_api.investigation.state_machines import apply_finding_decision
from flow_api.publishing.publication import PublicationService
from flow_api.publishing.renderers import render_html, render_pptx, render_xlsx
from flow_api.publishing.service import (
    PublishingFreezeError,
    freeze_report_snapshot,
)


def _approve_top_findings(session: Session, count: int = 1) -> None:
    findings = session.scalars(select(Finding).order_by(Finding.total_score.desc())).all()
    assert findings
    for finding in findings[:count]:
        evidence = session.scalars(select(Evidence).where(Evidence.finding_id == finding.id)).all()
        assert all(item.status == "verified" for item in evidence)
        session.add(
            Conclusion(
                finding=finding,
                verified_facts="已验证事实成立。",
                analysis_judgment="成本上升是主要驱动。",
                open_questions="暂无。",
                recommendation="重新议价。",
            )
        )
        session.flush()
        apply_finding_decision(session, finding, "submitted", reviewer="陈晨", comment=None)
        apply_finding_decision(session, finding, "approved", reviewer="王总", comment=None)


def test_freeze_requires_approved_findings(publishing_session: Session) -> None:
    publish_analysis_run(publishing_session)
    snapshot_id = publishing_session.scalar(select(Finding.metric_snapshot_id).limit(1))
    with pytest.raises(PublishingFreezeError) as error:
        freeze_report_snapshot(publishing_session, metric_snapshot_id=snapshot_id)
    assert "approved" in str(error.value)


def test_freeze_and_render_all_formats(publishing_session: Session) -> None:
    publish_analysis_run(publishing_session)
    _approve_top_findings(publishing_session)
    snapshot_id = publishing_session.scalar(select(Finding.metric_snapshot_id).limit(1))
    _report, view = freeze_report_snapshot(publishing_session, metric_snapshot_id=snapshot_id)

    assert view.findings, "approved findings must be frozen into the view"
    assert view.metrics, "metric values must be frozen into the view"
    key_values = view.key_values()
    assert view.identity.metric_snapshot_id == key_values["metric_snapshot_id"]

    presentation = Presentation(BytesIO(render_pptx(view)))
    pptx_text = "\n".join(
        shape.text_frame.text
        for slide in presentation.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert view.identity.title in pptx_text
    assert f"v{view.identity.report_version}" in pptx_text

    workbook = load_workbook(BytesIO(render_xlsx(view)))
    assert "指标结果" in workbook.sheetnames
    xlsx_values = {
        str(cell.value)
        for sheet in workbook.worksheets
        for row in sheet.iter_rows()
        for cell in row
        if cell.value is not None
    }
    assert key_values["metric_snapshot_id"] in xlsx_values

    html_document = render_html(view).decode("utf-8")
    assert view.identity.title in html_document
    assert "证据脚注" in html_document
    assert key_values["metric_snapshot_id"] in html_document


class FakeStore:
    """In-memory content-addressed store for fast, network-free tests."""

    def __init__(self) -> None:
        self.objects: dict[str, bytes] = {}

    def put_immutable(self, content: bytes, filename: str):
        import hashlib

        from flow_api.infrastructure.models.intake import StoredObject

        sha = hashlib.sha256(content).hexdigest()
        if sha not in self.objects:
            self.objects[sha] = content
        return StoredObject(
            sha256=sha,
            object_key=f"raw/{sha[:2]}/{sha}",
            size_bytes=len(content),
            content_type="application/octet-stream",
        )


def test_publication_attempts_persist_and_retry(publishing_session: Session) -> None:
    publish_analysis_run(publishing_session)
    _approve_top_findings(publishing_session)
    snapshot_id = publishing_session.scalar(select(Finding.metric_snapshot_id).limit(1))
    report, _view = freeze_report_snapshot(publishing_session, metric_snapshot_id=snapshot_id)

    def boom(_html: bytes) -> bytes:
        raise RuntimeError("printer offline")

    service = PublicationService(pdf_printer=boom, store=FakeStore())
    outcomes = service.publish(publishing_session, report.id)
    assert outcomes["pdf"] == "failed"
    assert outcomes["pptx"] == "succeeded"

    service_retry = PublicationService(
        pdf_printer=lambda _html: b"%PDF-1.4 fake", store=FakeStore()
    )
    outcomes_retry = service_retry.publish(publishing_session, report.id)
    assert outcomes_retry["pdf"] == "succeeded"

    attempt_formats = publishing_session.scalars(select(PublicationAttempt.format)).all()
    assert attempt_formats.count("pdf") >= 2


def test_succeeded_attempt_persists_reusable_stored_object(publishing_session: Session) -> None:
    """Task 6 契约：成功 attempt 必须持久化并复用 StoredObject 行（下载授权边界）。"""
    import hashlib

    from publishing.publishing_support import fresh_approved_report

    from flow_api.infrastructure.models.intake import StoredObject

    report, _view = fresh_approved_report(publishing_session)
    payload = b"%PDF-1.4 persisted object"
    service = PublicationService(pdf_printer=lambda _html: payload, store=FakeStore())
    outcomes = service.publish(publishing_session, report.id, formats=("pdf",))
    assert outcomes["pdf"] == "succeeded"

    attempt = publishing_session.scalar(
        select(PublicationAttempt).where(
            PublicationAttempt.report_snapshot_id == report.id,
            PublicationAttempt.format == "pdf",
            PublicationAttempt.status == "succeeded",
        )
    )
    assert attempt is not None
    assert attempt.stored_object_id is not None
    stored_row = publishing_session.get(StoredObject, attempt.stored_object_id)
    assert stored_row is not None
    assert stored_row.sha256 == hashlib.sha256(payload).hexdigest()
    assert stored_row.size_bytes == len(payload)
    assert stored_row.object_key.startswith("raw/")

    # 内容寻址：同 payload 再次发布复用同一 StoredObject 行
    service.publish(publishing_session, report.id, formats=("pdf",))
    duplicate_rows = publishing_session.scalars(
        select(StoredObject).where(StoredObject.sha256 == stored_row.sha256)
    ).all()
    assert len(duplicate_rows) == 1


def test_frozen_report_survives_review_changes_and_reloads(publishing_session: Session) -> None:
    from publishing.publishing_support import fresh_approved_report

    from flow_api.publishing.models import view_to_json
    from flow_api.publishing.service import build_report_view

    report, view = fresh_approved_report(publishing_session)
    original = view_to_json(view)
    finding = publishing_session.get(Finding, view.findings[0].finding_id)
    apply_finding_decision(publishing_session, finding, "returned", reviewer="review", comment=None)
    publishing_session.commit()
    publishing_session.expire_all()
    assert view_to_json(build_report_view(publishing_session, report)) == original
    assert render_html(build_report_view(publishing_session, report)) == render_html(view)


def test_freeze_versions_changed_content_and_reuses_identical_content(
    publishing_session: Session,
) -> None:
    from publishing.publishing_support import fresh_approved_report

    report, view = fresh_approved_report(publishing_session)
    same, _ = freeze_report_snapshot(
        publishing_session, metric_snapshot_id=report.metric_snapshot_id
    )
    assert same.id == report.id
    conclusion = publishing_session.scalar(
        select(Conclusion).where(Conclusion.finding_id == view.findings[0].finding_id)
    )
    conclusion.recommendation = "新的经批准建议"
    publishing_session.flush()
    changed, changed_view = freeze_report_snapshot(
        publishing_session, metric_snapshot_id=report.metric_snapshot_id
    )
    assert changed.id != report.id
    assert changed.version == report.version + 1
    assert changed_view.findings[0].conclusion["recommendation"] == "新的经批准建议"


def test_freeze_rejects_wrong_explicit_run(publishing_session: Session) -> None:
    from uuid import uuid4

    from publishing.publishing_support import fresh_approved_report

    report, _ = fresh_approved_report(publishing_session)
    with pytest.raises(PublishingFreezeError):
        freeze_report_snapshot(
            publishing_session,
            metric_snapshot_id=report.metric_snapshot_id,
            analysis_run_id=uuid4(),
        )


def test_legacy_report_cannot_rebuild_current_state(publishing_session: Session) -> None:
    from publishing.publishing_support import fresh_approved_report

    from flow_api.infrastructure.models.publishing import ReportSnapshot
    from flow_api.publishing.service import build_report_view

    report, _ = fresh_approved_report(publishing_session)
    legacy = ReportSnapshot(
        metric_snapshot_id=report.metric_snapshot_id,
        version=report.version + 1,
        title="legacy",
        template_code=report.template_code,
    )
    publishing_session.add(legacy)
    publishing_session.flush()
    with pytest.raises(PublishingFreezeError, match="重新冻结|frozen"):
        build_report_view(publishing_session, legacy)


def test_rejecting_evidence_returns_finding_and_preserves_frozen_history(
    publishing_session: Session,
) -> None:
    from publishing.publishing_support import fresh_approved_report

    from flow_api.infrastructure.models.analytics import ReviewEvent
    from flow_api.investigation.state_machines import apply_evidence_decision
    from flow_api.publishing.service import build_report_view

    report, view = fresh_approved_report(publishing_session)
    finding = publishing_session.get(Finding, view.findings[0].finding_id)
    evidence = publishing_session.scalar(select(Evidence).where(Evidence.finding_id == finding.id))
    apply_evidence_decision(
        publishing_session, evidence, "rejected", reviewer="review", comment="否决"
    )
    publishing_session.commit()
    assert finding.status == "in_review"
    events = publishing_session.scalars(
        select(ReviewEvent)
        .where(ReviewEvent.finding_id == finding.id)
        .order_by(ReviewEvent.sequence)
    ).all()
    assert {event.decision for event in events[-2:]} == {"returned", "evidence_rejected"}
    with pytest.raises(PublishingFreezeError):
        freeze_report_snapshot(publishing_session, metric_snapshot_id=report.metric_snapshot_id)
    assert build_report_view(publishing_session, report) == view


def test_freeze_checks_evidence_even_for_legacy_approved_finding(
    publishing_session: Session,
) -> None:
    from publishing.publishing_support import fresh_approved_report

    report, view = fresh_approved_report(publishing_session)
    evidence = publishing_session.scalar(
        select(Evidence).where(Evidence.finding_id == view.findings[0].finding_id)
    )
    evidence.status = "rejected"
    publishing_session.flush()
    with pytest.raises(PublishingFreezeError, match="evidence"):
        freeze_report_snapshot(publishing_session, metric_snapshot_id=report.metric_snapshot_id)


def test_new_analysis_run_does_not_rebind_old_report(publishing_session: Session) -> None:
    from pathlib import Path

    from publishing.publishing_support import (
        ANALYSIS_POLICY,
        approve_top_findings,
        fresh_approved_report,
    )

    from flow_api.analysis.policy import load_analysis_policy
    from flow_api.analysis.service import AnalysisRunService
    from flow_api.publishing.service import build_report_view

    report, original = fresh_approved_report(publishing_session)
    policy = load_analysis_policy(Path(ANALYSIS_POLICY))
    policy = policy.model_copy(
        update={
            "policy": policy.policy.model_copy(update={"engine_version": "review-test-v2"}),
            "policy_hash": "a" * 64,
        }
    )
    run = AnalysisRunService().create_run(
        publishing_session,
        snapshot_id=report.metric_snapshot_id,
        loaded_policy=policy,
    )
    approve_top_findings(publishing_session, run_id=run.id)
    assert build_report_view(publishing_session, report) == original
    new_report, new_view = freeze_report_snapshot(
        publishing_session,
        metric_snapshot_id=report.metric_snapshot_id,
        analysis_run_id=run.id,
    )
    assert new_report.version == report.version + 1
    assert new_view.identity.analysis_run_id == str(run.id)
    assert build_report_view(publishing_session, report) == original


def test_frozen_payload_is_immutable_in_database(publishing_session: Session) -> None:
    from publishing.publishing_support import fresh_approved_report
    from sqlalchemy import text
    from sqlalchemy.exc import DBAPIError

    report, _ = fresh_approved_report(publishing_session)
    with pytest.raises(DBAPIError, match="immutable"), publishing_session.begin_nested():
        publishing_session.execute(
            text("UPDATE report_snapshot SET frozen_view = '{}'::jsonb WHERE id = :id"),
            {"id": report.id},
        )


def test_freeze_waits_for_evidence_rejection_and_refreshes_approval(
    publishing_session: Session,
) -> None:
    from concurrent.futures import ThreadPoolExecutor, TimeoutError
    from threading import Event

    from publishing.publishing_support import fresh_approved_report

    from flow_api.investigation.state_machines import apply_evidence_decision

    report, view = fresh_approved_report(publishing_session)
    started = Event()
    with Session(publishing_session.bind) as writer, Session(publishing_session.bind) as reader:
        # Preload the old approval in the reader to exercise ORM identity-map refresh.
        stale = reader.get(Finding, view.findings[0].finding_id)
        assert stale.status == "approved"
        evidence = writer.scalar(select(Evidence).where(Evidence.finding_id == stale.id))
        apply_evidence_decision(writer, evidence, "rejected", reviewer="review", comment=None)

        def freeze_after_lock() -> None:
            started.set()
            try:
                freeze_report_snapshot(reader, metric_snapshot_id=report.metric_snapshot_id)
            finally:
                reader.rollback()

        with ThreadPoolExecutor(max_workers=1) as executor:
            pending = executor.submit(freeze_after_lock)
            assert started.wait(5)
            try:
                with pytest.raises(TimeoutError):
                    pending.result(timeout=0.2)
            finally:
                writer.commit()
            with pytest.raises(PublishingFreezeError, match="approved"):
                pending.result(timeout=5)
