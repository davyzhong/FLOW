from io import BytesIO

from pptx import Presentation

from flow_api.publishing.models import ReportView, SnapshotIdentity, SnapshotMetric
from flow_api.publishing.renderers import render_pptx


def sample_view(metric_count: int = 15) -> ReportView:
    return ReportView(
        identity=SnapshotIdentity(
            batch_id="batch",
            metric_snapshot_id="snapshot",
            analysis_run_id="run",
            report_snapshot_id="report",
            report_version=1,
            title="FLOW 月度财务分析",
            template_code="flow.monthly-review.v1",
            metric_engine_version="v1",
            analysis_engine_version="v1",
            generated_at="2026-09-04T12:00:00Z",
        ),
        metrics=tuple(
            SnapshotMetric(
                code=f"metric_{index}",
                name=f"财务指标 {index}",
                formula="sum(amount)",
                unit="CNY",
                definition_version=1,
                comparison="current",
                period=None,
                value="1234567.8900",
                budget="1200000.0000",
                variance="34567.8900",
            )
            for index in range(metric_count)
        ),
        findings=(),
        quality_summary={"blocking": 0},
        reconciliations=(),
    )


def test_pptx_text_has_real_bounds_and_overview_paginates() -> None:
    deck = Presentation(BytesIO(render_pptx(sample_view(35))))
    overview = [slide for slide in deck.slides if slide.shapes.title.text.startswith("经营概览")]
    for slide in deck.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame or not shape.text_frame.text:
                continue
            assert shape.width > 0 and shape.height > 0
            assert shape.left >= 0 and shape.top >= 0
            assert shape.left + shape.width <= deck.slide_width
            assert shape.top + shape.height <= deck.slide_height
    assert len(overview) >= 3
    text = "\n".join(
        shape.text for slide in overview for shape in slide.shapes if shape.has_text_frame
    )
    assert all(f"metric_{index} =" in text for index in range(35))
