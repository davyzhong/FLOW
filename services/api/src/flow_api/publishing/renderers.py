"""PPTX / XLSX / HTML renderers for a frozen Report Snapshot.

PDF printing is intentionally not implemented here: the frozen HTML is
printed by the publication gate through the repository's pinned Chromium
(`npx playwright pdf`), and the resulting bytes are attached to the
publication attempt by the caller.
"""

from __future__ import annotations

import html as html_module
from collections.abc import Callable

from flow_api.publishing.models import ReportView, SnapshotFinding


def _money(value: str) -> str:
    try:
        number = float(value)
    except ValueError:
        return value
    return f"{number / 10000:,.0f} 万"


def render_pptx(view: ReportView) -> bytes:
    from io import BytesIO

    from pptx import Presentation
    from pptx.util import Pt

    presentation = Presentation()
    identity = view.identity

    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = identity.title
    body = slide.placeholders[1].text_frame
    body.text = "结论先行：本报告全部数字来自已发布指标快照与已签发发现。"
    for line in [
        f"数据批次 {identity.batch_id} · 报告 v{identity.report_version}",
        f"指标快照 {identity.metric_snapshot_id}",
        f"分析运行 {identity.analysis_run_id}",
        f"指标引擎 {identity.metric_engine_version} · 分析引擎 {identity.analysis_engine_version}",
    ]:
        paragraph = body.add_paragraph()
        paragraph.text = line
        paragraph.font.size = Pt(12)

    overview = presentation.slides.add_slide(presentation.slide_layouts[5])
    overview.shapes.title.text = "经营概览"
    body_box = (
        overview.placeholders[1].text_frame
        if len(overview.placeholders) > 1
        else overview.shapes.add_textbox(0, 0, 0, 0).text_frame
    )
    body_box.text = f"共 {len(view.metrics)} 项指标、{len(view.findings)} 项已签发发现"
    for metric in view.metrics:
        paragraph = body_box.add_paragraph()
        paragraph.text = f"{metric.code} = {metric.value} {metric.unit}"
        paragraph.font.size = Pt(10)

    for metric in view.metrics[:8]:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = f"{metric.name}（{metric.code}）"
        frame = (
            slide.placeholders[1].text_frame
            if len(slide.placeholders) > 1
            else slide.shapes.add_textbox(0, 0, 0, 0).text_frame
        )
        frame.text = f"本期 {metric.value} {metric.unit}"
        if metric.budget:
            paragraph = frame.add_paragraph()
            paragraph.text = f"预算 {metric.budget} · 差异 {metric.variance or '—'}"

    if view.findings:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "已签发发现与建议"
        frame = slide.placeholders[1].text_frame
        frame.text = "结论 / 证据索引见附录"
        for finding in view.findings:
            paragraph = frame.add_paragraph()
            exact = finding.impact_amount
            paragraph.text = (
                f"{finding.title} · 影响 {_money(exact)}（精确 {exact}）"
            )

    evidence_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    evidence_slide.shapes.title.text = "证据索引"
    evidence_frame = (
        evidence_slide.placeholders[1].text_frame
        if len(evidence_slide.placeholders) > 1
        else evidence_slide.shapes.add_textbox(0, 0, 0, 0).text_frame
    )
    evidence_frame.text = "；".join(
        f"{finding.finding_id[:8]} ← {len(finding.evidence_ids)} 项证据"
        for finding in view.findings
    )

    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def render_xlsx(view: ReportView) -> bytes:
    from io import BytesIO

    from openpyxl import Workbook

    workbook = Workbook()
    metrics_sheet = workbook.active
    assert metrics_sheet is not None, "new workbook always has an active sheet"
    metrics_sheet.title = "指标结果"
    metrics_sheet.append(
        ["指标", "代码", "本期", "预算", "差异", "单位", "公式", "口径版本"]
    )
    for metric in view.metrics:
        metrics_sheet.append(
            [
                metric.name,
                metric.code,
                metric.value,
                metric.budget or "",
                metric.variance or "",
                metric.unit,
                metric.formula,
                metric.definition_version,
            ]
        )

    findings_sheet = workbook.create_sheet("已签发发现")
    findings_sheet.append(["发现ID", "标题", "类型", "影响金额", "状态"])
    for finding in view.findings:
        findings_sheet.append(
            [
                finding.finding_id,
                finding.title,
                finding.finding_type or "",
                finding.impact_amount,
                finding.status,
            ]
        )

    drivers_sheet = workbook.create_sheet("驱动明细")
    drivers_sheet.append(["发现ID", "驱动", "金额", "占比", "计算方法"])
    for finding in view.findings:
        for driver in finding.drivers:
            drivers_sheet.append(
                [
                    finding.finding_id,
                    driver["code"],
                    driver["amount"],
                    driver["ratio"],
                    driver["method"],
                ]
            )

    quality_sheet = workbook.create_sheet("质量与对账")
    quality_sheet.append(["检查", "结果"])
    for key, value in view.quality_summary.items():
        quality_sheet.append([f"{key} 问题数", value])
    for item in view.reconciliations:
        quality_sheet.append([f"对账 {item['code']}", item["passed"]])

    footer = workbook.create_sheet("版本与血缘")
    footer.append(["报告标题", view.identity.title])
    footer.append(["项目", "值"])
    for footer_key, footer_value in view.identity_footer().items():
        footer.append([footer_key, footer_value])

    buffer = BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def _escape(value: str) -> str:
    return html_module.escape(value, quote=True)


_STYLE = (
    "body{margin:2rem;color:#1b2a3f;"
    'font-family:"Inter Variable","Noto Sans SC Variable",sans-serif}'
    "table{border-collapse:collapse}"
    "td,th{border:1px solid #ccd6e0;padding:.35rem .6rem;font-size:.9rem}"
    ".finding{border-left:3px solid #49a389;padding-left:.8rem;margin:1rem 0}"
    "footer{margin-top:2rem;font-size:.8rem;color:#5c6b80}"
)


def render_html(view: ReportView) -> bytes:
    metric_rows = "".join(
        "<tr><th>{name}</th><td>{value}</td><td>{budget}</td><td>{variance}</td></tr>".format(
            name=_escape(metric.name),
            value=_escape(metric.value),
            budget=_escape(metric.budget or "—"),
            variance=_escape(metric.variance or "—"),
        )
        for metric in view.metrics
    )

    def _finding_section(finding: SnapshotFinding) -> str:
        driver_items = "".join(
            "<li>{} {}</li>".format(
                _escape(str(driver.get("code", ""))),
                _escape(_money(str(driver.get("amount", "")))),
            )
            for driver in finding.drivers
        )
        impact_html = _escape(_money(finding.impact_amount))
        impact_exact = _escape(finding.impact_amount)
        status_html = _escape(finding.status)
        head = (
            f"<p>影响金额 {impact_html} · 状态 {status_html}"
            f"<br><small>精确: {impact_exact}</small></p>"
        )
        return (
            '<section class="finding">'
            f"<h3>{_escape(finding.title)}</h3>"
            f"{head}"
            f"<ul>{driver_items}</ul>"
            "</section>"
        )

    finding_sections = "".join(_finding_section(finding) for finding in view.findings)

    footnotes = "".join(
        "<li>证据 {eid} · {fid}</li>".format(  # noqa: UP032
            eid=_escape(evidence_id), fid=_escape(finding.finding_id)
        )
        for finding in view.findings
        for evidence_id in finding.evidence_ids
    )
    footer_rows = "".join(
        f"<tr><th>{_escape(key)}</th><td>{_escape(value)}</td></tr>"
        for key, value in view.identity_footer().items()
    )
    metric_count = _escape(str(len(view.metrics)))
    document = f"""<!doctype html>
<html lang="zh-CN">
<head><meta charset="utf-8">
<title>{_escape(view.identity.title)}</title>
<style>{_STYLE}</style>
</head><body>
<h1>{_escape(view.identity.title)}</h1>
<p>结论先行：全部数字来自已发布指标快照与已签发发现，共 {metric_count} 项指标。</p>
<table><thead><tr><th>指标</th><th>本期</th><th>预算</th><th>差异</th></tr></thead><tbody>{metric_rows}</tbody></table>
<h2>已签发发现</h2>
{finding_sections or "<p>本报告未包含已签发发现。</p>"}
<h2>证据脚注</h2>
<ol>{footnotes or "<li>无</li>"}</ol>
<footer><table>{footer_rows}</table></footer>
</body></html>"""
    return document.encode("utf-8")


RENDERERS: dict[str, Callable[[ReportView], bytes]] = {
    "pptx": render_pptx,
    "xlsx": render_xlsx,
    "html": render_html,
}


__all__ = ["RENDERERS", "render_html", "render_pptx", "render_xlsx"]
